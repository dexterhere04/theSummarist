"""Text extraction from uploaded files and URLs."""
from __future__ import annotations

import io
from pathlib import Path
from urllib.parse import urlparse

from app.config import settings
from app.core.errors import ExtractionFailedError, InvalidURLError
from app.models.enums import DocFormat
from app.services.ocr import ocr_image

_BULLET_PREFIXES = ("- ", "• ", "* ")


def _split_paragraphs(text: str) -> list[str]:
    return [p.strip() for p in text.split("\n\n") if p.strip()]


def _build_body(title: str, paragraphs: list[str]) -> list[dict]:
    body_paragraphs: list[str] = []
    bullets: list[str] = []
    for para in paragraphs[:6]:
        if para.startswith(_BULLET_PREFIXES):
            bullets.append(para.lstrip("-•* ").strip())
        else:
            body_paragraphs.append(para)
    if not body_paragraphs:
        body_paragraphs = [title]
    section: dict = {"heading": title, "paragraphs": body_paragraphs}
    if bullets:
        section["bullets"] = bullets
    return [section]


def _result(
    title: str, body: list[dict], pages: int | None, words: int, ocr_method: str
) -> dict:
    return {
        "title": title,
        "body": body,
        "pages": pages,
        "words": words,
        "language": "English (US)",
        "ocr_method": ocr_method,
    }


async def extract_from_file(document, data: bytes) -> dict:
    fmt = document.format
    title = Path(document.file_name).stem

    try:
        if fmt == DocFormat.PDF:
            return await _extract_pdf(title, data)
        if fmt == DocFormat.DOCX:
            return await _extract_docx(title, data)
        if fmt == DocFormat.PPTX:
            return await _extract_pptx(title, data)
        if fmt == DocFormat.TXT:
            return _extract_txt(title, data)
        if fmt == DocFormat.DOC:
            raise ExtractionFailedError(
                "Legacy .doc extraction is not supported; convert to .docx or .pdf"
            )
        raise ExtractionFailedError(f"Unsupported format: {fmt.value}")
    except ExtractionFailedError:
        raise
    except Exception as e:
        raise ExtractionFailedError(f"Extraction failed: {e}") from e


async def _extract_pdf(title: str, data: bytes) -> dict:
    from pypdf import PdfReader

    reader = PdfReader(io.BytesIO(data))
    pages = len(reader.pages)
    text = "\n\n".join(page.extract_text() or "" for page in reader.pages)

    ocr_method = "Native Text Layer"
    if not text.strip():
        # Best-effort OCR fallback using the first embedded image on page one.
        ocr_method = settings.ocr_engine_label
        try:
            image = reader.pages[0].images[0]
            text = ocr_image(image.data)
        except (IndexError, KeyError):
            raise ExtractionFailedError(
                "This PDF has no extractable text layer and no embedded "
                "image to OCR"
            ) from None

    paragraphs = _split_paragraphs(text)
    words = len(text.split())
    return _result(title, _build_body(title, paragraphs), pages, words, ocr_method)


async def _extract_docx(title: str, data: bytes) -> dict:
    from docx import Document as DocxDocument

    doc = DocxDocument(io.BytesIO(data))
    text = "\n".join(p.text for p in doc.paragraphs)
    paragraphs = _split_paragraphs(text)
    words = len(text.split())
    return _result(title, _build_body(title, paragraphs), None, words, "Native Text Layer")


async def _extract_pptx(title: str, data: bytes) -> dict:
    from pptx import Presentation

    prs = Presentation(io.BytesIO(data))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    text = "\n".join(parts)
    paragraphs = _split_paragraphs(text)
    words = len(text.split())
    return _result(
        title, _build_body(title, paragraphs), len(prs.slides), words, "Native Text Layer"
    )


def _extract_txt(title: str, data: bytes) -> dict:
    try:
        text = data.decode("utf-8")
    except UnicodeDecodeError:
        text = data.decode("latin-1")
    paragraphs = _split_paragraphs(text)
    words = len(text.split())
    return _result(title, _build_body(title, paragraphs), 1, words, "Native Text Layer")


async def extract_from_url(url: str) -> dict:
    import trafilatura

    try:
        downloaded = trafilatura.fetch_url(url)
        if downloaded is None:
            raise InvalidURLError("Could not fetch URL")

        text = trafilatura.extract(
            downloaded, include_comments=False, include_tables=False
        )
        if not text:
            raise InvalidURLError("No extractable content")

        metadata = trafilatura.extract_metadata(downloaded)
        title = metadata.title if metadata and metadata.title else urlparse(url).netloc

        paragraphs = _split_paragraphs(text)
        body = [{"heading": title, "paragraphs": paragraphs}]
        words = len(text.split())
        return _result(title, body, 1, words, "Native Text Layer")
    except InvalidURLError:
        raise
    except Exception as e:
        raise InvalidURLError(f"Could not fetch URL: {e}") from e
